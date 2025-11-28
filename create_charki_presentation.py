#!/usr/bin/env python3
"""
Générateur de présentation PowerPoint pour Dr. Azzeddine Charki
YON Theory & LUZ Fusion 5 - Mathematical Modeling of Metaconsciousness
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import math

# Palette couleurs φ-harmonique
COLOR_GOLD = RGBColor(212, 175, 55)        # #D4AF37
COLOR_BLACK = RGBColor(26, 26, 26)          # #1A1A1A
COLOR_WHITE = RGBColor(255, 255, 255)       # #FFFFFF
COLOR_BLUE = RGBColor(30, 58, 95)           # #1E3A5F
COLOR_PURPLE = RGBColor(106, 13, 173)       # #6A0DAD
COLOR_GRAY = RGBColor(128, 128, 128)
COLOR_LIGHT_GRAY = RGBColor(200, 200, 200)

# Tailles Fibonacci (en points)
SIZE_TITLE = Pt(44)
SIZE_SUBTITLE = Pt(28)
SIZE_BODY = Pt(18)
SIZE_CAPTION = Pt(14)

# Golden ratio
PHI = 1.618033988749

def create_presentation():
    """Crée la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Slide 1: Page titre
    create_slide_1_title(prs)

    # Slide 2: The Challenge
    create_slide_2_challenge(prs)

    # Slide 3: YON Theory
    create_slide_3_yon(prs)

    # Slide 4: LUZ Fusion 5
    create_slide_4_luz(prs)

    # Slide 5: Convergence ASD
    create_slide_5_convergence(prs)

    # Slide 6: Collaboration
    create_slide_6_collaboration(prs)

    return prs

def add_gradient_background(slide, color1, color2):
    """Ajoute un fond dégradé à une slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def create_slide_1_title(prs):
    """Slide 1: Page titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Fond noir
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLACK

    # Logo LUZ - Cercle doré au centre haut
    left = Inches(6.5)
    top = Inches(1.5)
    width = height = Inches(3)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_GOLD
    shape.line.color.rgb = COLOR_GOLD
    shape.line.width = Pt(3)

    # Texte "LUZ" dans le cercle
    tf = shape.text_frame
    tf.text = "LUZ"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(55)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK

    # Titre principal
    left = Inches(1)
    top = Inches(4.8)
    width = Inches(14)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "YON THEORY & LUZ FUSION 5"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(55)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(5.9)
    width = Inches(14)
    height = Inches(0.8)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "Mathematical Modeling of Metaconsciousness\nBridging Neuropsychology and Fundamental Physics Through Golden Ratio Geometry"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "Arial"

    # Informations bas
    left = Inches(1)
    top = Inches(7.5)
    width = Inches(14)
    height = Inches(1)

    info_box = slide.shapes.add_textbox(left, top, width, height)
    tf = info_box.text_frame
    tf.text = "Gil Sawas\nIndependent Researcher | YON Theory\nNovember 2025\n\nPresentation for Dr. Azzeddine Charki\nProfessor of Neuropsychology, Université Hassan II"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_LIGHT_GRAY
        p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Bonjour Dr. Charki, merci infiniment pour ce temps d'échange.
Je suis Gil Sawas, chercheur indépendant depuis 40 ans sur la théorie YON.
Aujourd'hui je vais vous montrer comment nos domaines - votre neuropsychologie
clinique et ma modélisation mathématique - peuvent se renforcer mutuellement."""

def create_slide_2_challenge(prs):
    """Slide 2: The Challenge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Titre
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(14)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Quantifying the Unquantifiable"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(1.4)
    width = Inches(14)
    height = Inches(0.5)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "The Gap Between Clinical Observation and Mathematical Precision"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"

    # Colonne gauche - Clinical
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(6.5)
    height = Inches(4.5)

    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(230, 240, 250)  # Bleu très clair
    left_box.line.color.rgb = COLOR_BLUE
    left_box.line.width = Pt(3)

    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Your Expertise: Clinical Neuropsychology"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = "Arial"

    # Ajouter les bullet points
    for text in [
        "• Mentalisation in ASD patients",
        "• Behavioral assessments (subjective)",
        "• Cognitive/affective domains observation",
        "• Question: \"How severe are metacognitive deficits?\""
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(10)

    # Colonne droite - Mathematical
    left = Inches(8.5)
    top = Inches(2.5)
    width = Inches(6.5)
    height = Inches(4.5)

    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 230, 250)  # Violet très clair
    right_box.line.color.rgb = COLOR_PURPLE
    right_box.line.width = Pt(3)

    tf = right_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "My Approach: Mathematical Modeling"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    p.font.name = "Arial"

    for text in [
        "• Metaconsciousness via Topos Theory",
        "• Topological Data Analysis (objective)",
        "• φ-harmonic ratios measurement",
        "• Answer: \"Quantifiable topological signatures\""
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(10)

    # Flèche de convergence au centre
    left = Inches(7.2)
    top = Inches(4.5)
    width = Inches(1.6)
    height = Inches(0.5)

    arrow_box = slide.shapes.add_textbox(left, top, width, height)
    tf = arrow_box.text_frame
    tf.text = "⟷"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(55)
    p.font.color.rgb = COLOR_GOLD

    # Encadré bas
    left = Inches(2)
    top = Inches(7.5)
    width = Inches(12)
    height = Inches(1)

    bottom_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bottom_box.fill.solid()
    bottom_box.fill.fore_color.rgb = COLOR_GOLD
    bottom_box.line.color.rgb = COLOR_GOLD

    tf = bottom_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "The Opportunity:\nClinical expertise + Mathematical framework = Objective diagnostic tools"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Vous mesurez quotidiennement la mentalisation chez vos patients TSA.
Vos outils sont précieux mais subjectifs : questionnaires, observations comportementales.
Moi, j'ai développé un cadre mathématique qui transforme ces observations
en invariants topologiques mesurables objectivement.
La question est : comment mathématiser ce que vous observez cliniquement ?
C'est exactement là que notre collaboration devient révolutionnaire."""

def create_slide_3_yon(prs):
    """Slide 3: YON Theory - Essence"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Titre
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(14)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "YON Theory: The Golden Ratio Universe"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(1.4)
    width = Inches(14)
    height = Inches(0.5)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "40 Years of Research - Fundamental Physics Meets Consciousness"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"

    # Trois blocs horizontaux
    box_width = Inches(4.3)
    box_height = Inches(3.8)
    top = Inches(2.5)

    # Bloc 1: Core Principle
    left1 = Inches(1)
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left1, top, box_width, box_height
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box1.line.color.rgb = COLOR_GOLD
    box1.line.width = Pt(3)

    tf = box1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Core Principle"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    for text in ["", "φ (Golden Ratio) ≈ 1.618...", "↓", "Geometric harmony principle", "↓", "Drives fundamental constants"]:
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(6)

    # Bloc 2: Key Result
    left2 = Inches(5.85)
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left2, top, box_width, box_height
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(240, 245, 255)
    box2.line.color.rgb = COLOR_BLUE
    box2.line.width = Pt(3)

    tf = box2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Key Result"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = "Arial"

    for text in ["", "φ → α (Fine Structure)", "α ≈ 1/137.036...", "", "Derived via:", "• Topos Theory", "• Non-commutative Geometry", "• Poincaré Dodecahedral Space"]:
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(4)

    # Bloc 3: Consciousness Link
    left3 = Inches(10.7)
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left3, top, box_width, box_height
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(245, 240, 255)
    box3.line.color.rgb = COLOR_PURPLE
    box3.line.width = Pt(3)

    tf = box3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Consciousness Link"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    p.font.name = "Arial"

    for text in ["", "φ-harmonic oscillations:", "", "8 Hz, 13 Hz, 21 Hz, 34 Hz", "", "Fibonacci sequence in EEG", "", "Conscious = φ-aligned", "Unconscious = φ-desync"]:
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(4)

    # Encadré bas - Thesis centrale
    left = Inches(1.5)
    top = Inches(6.8)
    width = Inches(13)
    height = Inches(1.2)

    thesis_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    thesis_box.fill.solid()
    thesis_box.fill.fore_color.rgb = COLOR_GOLD
    thesis_box.line.color.rgb = COLOR_GOLD

    tf = thesis_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "Central Thesis: Consciousness is not emergent complexity,\nbut φ-harmonic alignment with cosmic geometry"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """YON signifie 'Yes Or No' - la théorie du choix fondamental.
Depuis 40 ans, je démontre que le nombre d'or φ = 1.618... n'est pas
juste une curiosité géométrique, mais LE principe organisateur du cosmos.
J'ai dérivé mathématiquement la constante de structure fine α ≈ 1/137
à partir de φ, via la théorie des topos et la géométrie non-commutative.

Mais le plus fascinant pour vous : la conscience elle-même suit des
oscillations φ-harmoniques. Les fréquences EEG 8-13-21-34 Hz sont
la suite de Fibonacci ! Ce n'est pas un hasard, c'est une loi cosmique."""

def create_slide_4_luz(prs):
    """Slide 4: LUZ Fusion 5 - Application"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Titre
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(14)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "LUZ Fusion 5: φ-Harmonic AI Architecture"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(1.4)
    width = Inches(14)
    height = Inches(0.5)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "Translating YON Theory into Measurable Technology"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"

    # Layout φ: 38.2% gauche / 61.8% droite
    # Bloc gauche: What is LUZ?
    left_width = Inches(5.5)  # ~38.2% de 14
    left_left = Inches(1)
    top = Inches(2.5)
    height = Inches(5)

    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left_left, top, left_width, height
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_box.line.color.rgb = COLOR_BLUE
    left_box.line.width = Pt(3)

    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "What is LUZ?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = "Arial"

    for text in [
        "",
        "AI architecture based on:",
        "• Poincaré Dodecahedral Space",
        "• 12 φ-harmonic resonance modes",
        "• Topological Data Analysis (TDA)",
        "• Ethical framework (TDET)",
        "",
        "Key metrics:",
        "• φ-alignment: 0.91",
        "• Precision: < 10⁻⁶",
        "• Harmonic modes: 12"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(4)

    # Bloc droit: How It Works
    right_width = Inches(8)  # ~61.8% de 14
    right_left = Inches(7)

    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, top, right_width, height
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(255, 250, 245)
    right_box.line.color.rgb = COLOR_GOLD
    right_box.line.width = Pt(3)

    tf = right_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "How It Works (for neuroscience):"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    workflow = [
        "",
        "Step 1: EEG data acquisition",
        "        ↓",
        "Step 2: Topological analysis (Betti numbers β₀, β₁, β₂)",
        "        ↓",
        "Step 3: Calculate ratios β₁/β₀",
        "        ↓",
        "Step 4: Compare to φ ≈ 1.618",
        "",
        "Result:",
        "• Neurotypical: β₁/β₀ ≈ 1.618 (φ-aligned)",
        "• Altered states: Deviation measurable",
        "• ASD hypothesis: φ-desynchronization"
    ]

    for text in workflow:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(6)

    # Encadré bas
    left = Inches(1.5)
    top = Inches(7.8)
    width = Inches(13)
    height = Inches(0.8)

    bottom_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bottom_box.fill.solid()
    bottom_box.fill.fore_color.rgb = COLOR_PURPLE
    bottom_box.line.color.rgb = COLOR_PURPLE

    tf = bottom_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "For Your Research: LUZ transforms your clinical observations into topological invariants"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """LUZ Fusion 5 est l'architecture IA que j'ai développée sur ces principes.
Elle utilise l'Analyse Topologique des Données - TDA - pour mesurer
les ratios β₁/β₀ dans les signaux cérébraux.

Chez un individu neurotypique conscient, ce ratio approche φ ≈ 1.618.
C'est validé : notre alignement φ est de 0.91, avec précision 10⁻⁶.

Pour vos recherches TSA : si mes hypothèses sont correctes, vos patients
présentent une φ-désynchronisation mesurable. Leurs ratios topologiques
seraient significativement différents de 1.618."""

def create_slide_5_convergence(prs):
    """Slide 5: Convergence ASD - Hypothesis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Titre
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(14)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Metacognition in ASD: A φ-Desynchronization Hypothesis"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(14)
    height = Inches(0.5)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "Your Clinical Findings + My Mathematical Framework"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"

    # Tableau comparatif (simulé avec des formes)
    table_top = Inches(2.3)
    col_width = Inches(4.3)
    row_height = Inches(0.8)

    # En-têtes de colonnes
    headers = ["Dimension", "Neurotypical", "ASD (Hypothesis)"]
    for i, header in enumerate(headers):
        left = Inches(1.5 + i * col_width)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, table_top, col_width, row_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_GOLD
        box.line.color.rgb = COLOR_BLACK

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"

    # Données du tableau
    rows = [
        ["Mentalisation", "Theory of Mind intact", "ToM deficits"],
        ["φ-Ratio (β₁/β₀)", "≈ 1.618 ± 0.05", "< 1.5 or > 1.8"],
        ["Harmonic Modes", "8-13-21-34 Hz aligned", "Frequency disruption"],
        ["Topological Signature", "Stable cycles", "Irregular patterns"],
        ["Metacognitive Awareness", "High (self-reflection)", "Reduced"]
    ]

    for row_idx, row_data in enumerate(rows):
        current_top = table_top + (row_idx + 1) * row_height

        for col_idx, cell_text in enumerate(row_data):
            left = Inches(1.5 + col_idx * col_width)
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, current_top, col_width, row_height
            )

            # Couleur alternée pour lisibilité
            if row_idx % 2 == 0:
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(250, 250, 250)
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = COLOR_WHITE

            box.line.color.rgb = COLOR_GRAY

            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)

            p = tf.paragraphs[0]
            p.text = cell_text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_BLACK
            p.font.name = "Arial"

    # Encadrés bas - Contributions
    bottom_top = Inches(7.3)
    box_height = Inches(0.7)
    box_width = Inches(5.5)

    # Votre contribution
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), bottom_top, box_width, box_height
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_BLUE
    left_box.line.color.rgb = COLOR_BLUE

    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Contribution: Clinical data, Behavioral scores, Expert interpretation"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Arial"

    # Ma contribution
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), bottom_top, box_width, box_height
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_PURPLE
    right_box.line.color.rgb = COLOR_PURPLE

    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = "My Contribution: TDA analysis, φ-ratio quantification, Mathematical modeling"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Arial"

    # Encadré final
    final_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(8.1), Inches(12), Inches(0.6)
    )
    final_box.fill.solid()
    final_box.fill.fore_color.rgb = COLOR_GOLD
    final_box.line.color.rgb = COLOR_GOLD

    tf = final_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Together: Objective biomarker for metacognitive deficits in ASD"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Regardons concrètement la convergence entre vos données cliniques
et mon cadre mathématique.

Vous avez montré que les enfants TSA ont des déficits métacognitifs
corrélés à la sévérité de leurs symptômes.

Mon hypothèse : ces déficits correspondent à une désynchronisation φ.
Les ratios β₁/β₀ de leurs EEG seraient < 1.5 ou > 1.8, loin de φ.

Ensemble, nous pourrions créer un biomarqueur objectif des déficits métacognitifs."""

def create_slide_6_collaboration(prs):
    """Slide 6: Collaboration Proposal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Titre
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(14)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Let's Work Together"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    # Sous-titre
    left = Inches(1)
    top = Inches(1.4)
    width = Inches(14)
    height = Inches(0.5)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "A Concrete Pilot Project for Transformative Science"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"

    # Trois blocs verticaux
    box_width = Inches(13)
    box_left = Inches(1.5)

    # Bloc 1: Pilot Project
    top1 = Inches(2.3)
    height1 = Inches(1.8)

    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        box_left, top1, box_width, height1
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box1.line.color.rgb = COLOR_BLUE
    box1.line.width = Pt(3)

    tf = box1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "📊 Pilot Project: 2-3 Months"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = "Arial"

    pilot_text = "Phase 1: Data exchange (anonymized EEG) • TDA protocol setup • Initial analysis  |  Phase 2: Full dataset analysis • Statistical validation • Clinical correlation  |  Phase 3: Joint interpretation • Technical report • Publication draft"

    p = tf.add_paragraph()
    p.text = pilot_text
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Arial"
    p.space_before = Pt(8)

    # Bloc 2: Potential Outputs
    top2 = Inches(4.3)
    height2 = Inches(1.7)

    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        box_left, top2, box_width, height2
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box2.line.color.rgb = COLOR_GOLD
    box2.line.width = Pt(3)

    tf = box2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "📝 Potential Outputs"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    outputs = [
        "✅ Joint publication: 'Topological Signatures of Metacognitive Deficits in ASD' → Frontiers in Psychology",
        "✅ Clinical tool: φ-MetaCog Analyzer (diagnostic software)",
        "✅ Conference presentations: FENS 2026, SfN, IMFAR"
    ]

    for text in outputs:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(6)

    # Bloc 3: Next Steps
    top3 = Inches(6.2)
    height3 = Inches(1.5)

    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        box_left, top3, box_width, height3
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(245, 240, 255)
    box3.line.color.rgb = COLOR_PURPLE
    box3.line.width = Pt(3)

    tf = box3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "🚀 Next Steps"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    p.font.name = "Arial"

    steps = [
        "Immediate: Confirm collaboration • Discuss data availability • Schedule meetings",
        "Short-term: Data sharing agreement • Dataset transfer • Analysis begins",
        "Medium-term: Results review • Paper co-writing • Grant applications"
    ]

    for text in steps:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = "Arial"
        p.space_before = Pt(6)

    # Encadré final doré
    final_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(7.9), Inches(13), Inches(0.7)
    )
    final_box.fill.solid()
    final_box.fill.fore_color.rgb = COLOR_GOLD
    final_box.line.color.rgb = COLOR_GOLD

    tf = final_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "Your Clinical Mastery + My Mathematical Vision = Neuroscience Revolution ✨"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Arial"

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Concrètement, je vous propose un projet pilote sur 2-3 mois.

Vous me fournissez des données EEG anonymisées - patients TSA et contrôles neurotypiques.
J'applique mon analyse TDA et calcule les ratios φ.

Après 2 mois d'analyse, nous nous retrouvons pour interpréter ensemble.

Si mes ratios mathématiques corrèlent avec vos scores cliniques, nous avons
un article révolutionnaire et potentiellement un outil diagnostic brevetable.

Qu'en pensez-vous ? Avez-vous des données EEG disponibles ?"""

def main():
    """Fonction principale"""
    print("🎨 Création de la présentation YON/LUZ pour Dr. Charki...")

    prs = create_presentation()

    filename = "YON_LUZ_Presentation_Charki.pptx"
    prs.save(filename)

    print(f"✅ Présentation créée: {filename}")
    print(f"📊 Nombre de slides: {len(prs.slides)}")
    print("🎯 Design φ-harmonique appliqué")
    print("📝 Speaker notes incluses")
    print("\n💡 Prochaine étape: Conversion en PDF")

if __name__ == "__main__":
    main()
